"""PowerPoint generation service using python-pptx."""

import re
from io import BytesIO
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# Brand colors (as RGB tuples)
COLOR_BACKGROUND = RGBColor(0xFA, 0xF9, 0xF7)  # Warm off-white
COLOR_SURFACE = RGBColor(0xFF, 0xFF, 0xFF)  # White
COLOR_TEXT_PRIMARY = RGBColor(0x2D, 0x2A, 0x26)  # Warm charcoal
COLOR_TEXT_SECONDARY = RGBColor(0x6B, 0x65, 0x60)  # Warm gray
COLOR_ACCENT = RGBColor(0x7C, 0x98, 0x85)  # Soft sage green
COLOR_ACCENT_WARM = RGBColor(0xC4, 0xA7, 0x7D)  # Warm gold/tan
COLOR_HIGHLIGHT = RGBColor(0xF5, 0xF2, 0xED)  # Linen-like


def _parse_markdown_sections(content: str) -> list[dict]:
    """
    Parse markdown into sections based on H2 headers.

    Returns list of dicts with 'title' and 'content' keys.
    """
    sections = []
    current_section = None

    lines = content.split("\n")

    for line in lines:
        # H1 header - skip (used for cover)
        if line.startswith("# ") and not line.startswith("## "):
            continue

        # H2 header - new section
        if line.startswith("## "):
            if current_section:
                sections.append(current_section)
            current_section = {
                "title": line[3:].strip(),
                "content": [],
            }
        elif current_section is not None:
            current_section["content"].append(line)

    # Add last section
    if current_section:
        sections.append(current_section)

    # Join content back together
    for section in sections:
        section["content"] = "\n".join(section["content"]).strip()

    return sections


def _extract_table(content: str) -> list[list[str]] | None:
    """Extract first markdown table from content."""
    lines = content.split("\n")
    table_lines = []
    in_table = False

    for line in lines:
        if "|" in line:
            # Skip separator rows
            if re.match(r"^\|[\s\-:|]+\|$", line.strip()):
                continue
            in_table = True
            # Parse table row
            cells = [cell.strip() for cell in line.strip().split("|")]
            # Remove empty cells from edges
            cells = [c for c in cells if c]
            if cells:
                table_lines.append(cells)
        elif in_table:
            break

    return table_lines if len(table_lines) > 1 else None


def _extract_bullet_points(content: str) -> list[str]:
    """Extract bullet points from markdown content."""
    bullets = []
    lines = content.split("\n")

    for line in lines:
        stripped = line.strip()
        # Match bullet points and numbered lists
        if stripped.startswith("- ") or stripped.startswith("* "):
            text = stripped[2:].strip()
            # Remove markdown bold
            text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
            bullets.append(text)
        elif re.match(r"^\d+\.\s+", stripped):
            text = re.sub(r"^\d+\.\s+", "", stripped)
            text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
            bullets.append(text)

    return bullets


def _set_slide_background(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_slide(prs: Presentation, customer_name: str):
    """Add title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, COLOR_BACKGROUND)

    # Logo placeholder
    logo = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5),
        Inches(1.5),
        Inches(0.8),
        Inches(0.8),
    )
    logo.fill.solid()
    logo.fill.fore_color.rgb = COLOR_ACCENT
    logo.line.fill.background()

    # Add "ARI" text to logo
    text_frame = logo.text_frame
    text_frame.word_wrap = False
    p = text_frame.paragraphs[0]
    p.text = "ARI"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI Visibility Analysis"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = COLOR_TEXT_PRIMARY

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(9), Inches(0.6))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = customer_name
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_ACCENT

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.4))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Prepared by NewsUSA AI Intelligence Division | {datetime.now().strftime('%B %Y')}"
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY


def _add_section_slide(prs: Presentation, title: str, content: str):
    """Add a content slide with title and body."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, COLOR_BACKGROUND)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_TEXT_PRIMARY

    # Accent line under title
    line = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5),
        Inches(1.1),
        Inches(2),
        Inches(0.05),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    # Check for table in content
    table_data = _extract_table(content)
    if table_data:
        _add_table_to_slide(slide, table_data)
        return

    # Check for bullet points
    bullets = _extract_bullet_points(content)
    if bullets:
        _add_bullets_to_slide(slide, bullets)
        return

    # Plain text content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    # Clean up markdown formatting
    clean_content = content
    clean_content = re.sub(r"\*\*([^*]+)\*\*", r"\1", clean_content)  # Bold
    clean_content = re.sub(r"\*([^*]+)\*", r"\1", clean_content)  # Italic
    clean_content = re.sub(r"^---+$", "", clean_content, flags=re.MULTILINE)  # HR

    # Add text
    for i, para in enumerate(clean_content.split("\n\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = para.strip()
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT_PRIMARY
        p.space_after = Pt(12)


def _add_table_to_slide(slide, table_data: list[list[str]]):
    """Add a table to the slide."""
    rows = len(table_data)
    cols = len(table_data[0]) if table_data else 0

    if rows == 0 or cols == 0:
        return

    # Create table
    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(0.4 * rows),
    ).table

    # Style table
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)

            # Clean markdown from cell text
            text = re.sub(r"\*\*([^*]+)\*\*", r"\1", cell_text)
            cell.text = text

            # Style cell
            cell.text_frame.paragraphs[0].font.size = Pt(11)

            if row_idx == 0:
                # Header row
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT_PRIMARY
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_HIGHLIGHT
            else:
                cell.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT_PRIMARY


def _add_bullets_to_slide(slide, bullets: list[str]):
    """Add bullet points to the slide."""
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets[:8]):  # Max 8 bullets per slide
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT_PRIMARY
        p.space_after = Pt(10)
        p.level = 0


def _add_closing_slide(prs: Presentation, customer_name: str):
    """Add closing/contact slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, COLOR_BACKGROUND)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Next Steps"
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = COLOR_TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(2))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Ready to improve {customer_name}'s AI visibility?"
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p.alignment = PP_ALIGN.CENTER

    # Contact
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.5))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Contact NewsUSA AI Intelligence Division"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER


def generate_pptx(markdown_content: str, customer_name: str) -> bytes:
    """
    Generate a branded PowerPoint from client proposal markdown.

    Args:
        markdown_content: Raw markdown string (client proposal)
        customer_name: Customer display name

    Returns:
        PPTX as bytes
    """
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add title slide
    _add_title_slide(prs, customer_name)

    # Parse sections
    sections = _parse_markdown_sections(markdown_content)

    # Key sections to include (prioritize important content)
    priority_sections = [
        "The Opportunity",
        "Key Findings",
        "Why This Is Happening",
        "The Solution",
        "Recommended Articles",
        "Expected Outcomes",
        "Investment & Timeline",
        "Next Steps",
    ]

    # Add priority sections first
    added_sections = set()
    for priority_title in priority_sections:
        for section in sections:
            if priority_title.lower() in section["title"].lower():
                _add_section_slide(prs, section["title"], section["content"])
                added_sections.add(section["title"])
                break

    # Add remaining sections (up to a reasonable limit)
    for section in sections:
        if section["title"] not in added_sections and len(prs.slides) < 15:
            # Skip appendix-like sections
            if "appendix" in section["title"].lower():
                continue
            _add_section_slide(prs, section["title"], section["content"])

    # Add closing slide
    _add_closing_slide(prs, customer_name)

    # Save to bytes
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return buffer.read()
